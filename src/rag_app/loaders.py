"""Document loaders for PDF, DOCX, PPTX, XLSX.

Each loader returns the raw text of the document. A `Document` bundle
also carries metadata (source filename, format, domain) that is later
propagated to each chunk for citation in the final answer.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Dict, List, Union

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class Document:
    """A loaded document with provenance metadata."""
    text: str
    source: str
    format: str
    domain: str


# --- Per-format loaders ------------------------------------------------------

def load_pdf(path: str) -> str:
    """Extract text from a PDF, page by page."""
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_docx(path: str) -> str:
    """Extract paragraphs + table cells from a .docx file."""
    doc = DocxDocument(path)
    parts: List[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def load_pptx(path: str) -> str:
    """Extract text from each slide of a .pptx, prefixing slide index."""
    prs = Presentation(path)
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: List[str] = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if len(slide_text) > 1:
            parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def load_xlsx(path: str) -> str:
    """Render every sheet of an .xlsx as text (sheet name header + tabular string)."""
    sheets = pd.read_excel(path, sheet_name=None)
    parts: List[str] = []
    for sheet_name, df in sheets.items():
        parts.append(f"[Foglio: {sheet_name}]")
        parts.append(df.to_string(index=False))
    return "\n\n".join(parts)


LOADERS: Dict[str, Callable[[str], str]] = {
    ".pdf":  load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


# --- Public API --------------------------------------------------------------

def load_document(path: Union[str, Path], domain: str = "default") -> Document:
    """Load a single file and wrap it in a Document with metadata.

    Raises:
        ValueError: if the file extension is not supported.
    """
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in LOADERS:
        raise ValueError(
            f"Unsupported format: {ext}. Supported: {list(LOADERS)}"
        )
    text = LOADERS[ext](str(path))
    return Document(text=text, source=path.name, format=ext, domain=domain)


def load_directory(
    data_dir: Union[str, Path],
    domain_map: Dict[str, str],
    verbose: bool = False,
) -> List[Document]:
    """Load every supported file under each domain subfolder.

    Args:
        data_dir: root folder that contains the subfolders named in `domain_map`.
        domain_map: mapping `subfolder_name -> domain_label`.
        verbose: print a line per loaded file.

    Returns:
        List of Document objects. Unsupported extensions are skipped silently.
    """
    data_dir = Path(data_dir)
    docs: List[Document] = []
    for subfolder, domain in domain_map.items():
        folder = data_dir / subfolder
        if not folder.exists():
            continue
        for path in sorted(folder.iterdir()):
            if path.suffix.lower() not in LOADERS:
                continue
            try:
                doc = load_document(path, domain=domain)
                docs.append(doc)
                if verbose:
                    print(f"{path.name:55s} | {domain:10s} | {len(doc.text):>6d} char")
            except Exception as e:  # noqa: BLE001 — report and skip bad files
                print(f"[WARN] Could not load {path.name}: {e}")
    return docs
